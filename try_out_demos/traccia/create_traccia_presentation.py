from pptx import Presentation
from pptx.util import Inches, Pt

# Load the data summary content (manually copied from earlier read)
summary = {
    "intro": [
        "Founded in 2016, Traccia specializes in training programs at the nexus of business and IT.",
        "Small public company headquartered in Laren NH, Netherlands.",
        "Focus on enterprise architecture, requirements engineering, process design, and change management.",
    ],
    "history": [
        "Traccia develops training and masterclasses based on consulting success from Acumen.",
        "Offers Business Engineering course and masterclasses for professionals bridging business and IT.",
    ],
    "offerings": [
        "Training programs and masterclasses teaching advanced methods and best practices.",
        "Courses include Business Engineering and tailored masterclasses.",
    ],
    "help": [
        "Provides structured knowledge and practical skills to align business goals with IT architecture.",
        "Enhances requirements engineering, process design, and organizational change management.",
        "Supports enterprise architecture and IT development alignment.",
    ],
    "market": [
        "Niche focus on business-IT integration training with proven consulting foundations.",
        "Small team, positioned to scale amid digital transformation trends.",
        "Opportunity to meet growing demand for business-IT and organizational transformation training.",
    ],
    "conclusion": [
        "Traccia offers a unique investment opportunity in a growing market niche.",
        "Established curriculum with demonstrated success and practical relevance.",
        "Ideal for investors seeking scale in business-IT training amid digital transformation.",
    ]
}

# Create presentation
prs = Presentation()

# Helper function to add a slide with title and bullet points

def add_bullet_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear any existing content
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    return slide


# Slides content
slides = [
    ("Introduction to Traccia", summary["intro"]),
    ("What is Traccia and Its History", summary["history"]),
    ("Traccia's Training Programs and Course Offerings", summary["offerings"]),
    ("How Traccia Helps Businesses and Professionals", summary["help"]),
    ("Market Position and Opportunities", summary["market"]),
    ("Conclusion and Call to Action for Investors", summary["conclusion"]),
]

# Create slides
for title, content in slides:
    add_bullet_slide(prs, title, content)

# Save file
output_path = r"C:/Users/Dancy Naik/Documents/VS_Code_Test/wow_ai/try_out_demos/traccia/Traccia_Investor_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
